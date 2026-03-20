# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Define constants for slide layout indices
BLANK_SLIDE_LAYOUT_INDEX = 6

class OAISPresentation:
    """
    A class to generate PowerPoint presentations with a predefined OAIS style.

    This class provides methods to add various types of slides (cover, about, grid,
    features, tips, closing) and handles common formatting tasks like setting
    backgrounds, adding textboxes, and managing colors.
    """
    def __init__(self, theme='dark'):
        """
        Initializes the OAISPresentation with a given theme.

        Args:
            theme (str): The visual theme for the presentation ('dark' or 'light').
                         Currently, only 'dark' is fully supported for colors.
        """
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(10)
        self.presentation.slide_height = Inches(5.625)
        self.theme = theme
        
        # Define color palette based on theme
        self.colors = {
            'dark_bg': RGBColor(15, 15, 15),
            'card_bg': RGBColor(26, 26, 26),
            'light_bg': RGBColor(245, 245, 240),
            'white_text': RGBColor(255, 255, 255),
            'dark_text': RGBColor(26, 26, 26),
            'grey_text1': RGBColor(136, 136, 136),
            'grey_text2': RGBColor(102, 102, 102),
            'grey_text3': RGBColor(85, 85, 85),
            'mute_text': RGBColor(68, 68, 68),
        }

    def _add_slide(self):
        """
        Adds a new blank slide to the presentation.

        Returns:
            Slide: The newly added slide object.
        """
        slide_layout = self.presentation.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX] # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        return slide

    def _set_background(self, slide, color):
        """
        Sets the background color of a given slide.

        Args:
            slide (Slide): The slide object to modify.
            color (RGBColor): The RGB color to set as the background.
        """
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_cover_slide(self, title, subtitle=None, description=None, badge_text=None, footer_items=None):
        """
        Adds a cover slide to the presentation.

        Args:
            title (str): The main title of the cover slide.
            subtitle (str, optional): A subtitle for the cover slide. Defaults to None.
            description (str, optional): A brief description. Defaults to None.
            badge_text (str, optional): Text for a badge/label. Defaults to None.
            footer_items (list, optional): A list of dictionaries for footer information
                                          (e.g., [{"label": "Date", "value": "2025"}]).
                                          Defaults to None.
        """
        slide = self._add_slide()
        self._set_background(slide, self.colors['dark_bg'])

        if badge_text:
            self._add_textbox(slide, badge_text, Inches(0.5), Inches(0.5), Inches(2), Inches(0.5), font_size=Pt(10), font_color=self.colors['grey_text1'])

        self._add_textbox(slide, title, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5), font_size=Pt(44), font_bold=True, font_color=self.colors['white_text'])

        if subtitle:
            self._add_textbox(slide, subtitle, Inches(0.5), Inches(3.0), Inches(9), Inches(0.75), font_size=Pt(24), font_color=self.colors['grey_text1'])
        
        if description:
            self._add_textbox(slide, description, Inches(0.5), Inches(3.75), Inches(9), Inches(0.5), font_size=Pt(14), font_color=self.colors['grey_text2'])

        if footer_items:
            for i, item in enumerate(footer_items):
                self._add_textbox(slide, f"{item['label']}: {item['value']}", Inches(0.5 + i*2.5), Inches(4.8), Inches(2), Inches(0.5), font_size=Pt(10), font_color=self.colors['grey_text3'])
    
    def add_about_slide(self, icon, name, headline, subtitle=None, badge_text=None, description=None, stats=None):
        """
        Adds an "About" slide to the presentation, typically for an overview or summary.

        Args:
            icon (str): A character or short string for an icon/symbol.
            name (str): The main name or topic.
            headline (str): A prominent headline text.
            subtitle (str, optional): A subtitle. Defaults to None.
            badge_text (str, optional): Text for a badge/label. Defaults to None.
            description (str, optional): A longer description. Defaults to None.
            stats (list, optional): A list of dictionaries for statistics
                                   (e.g., [{"value": "71%", "label": "Coverage"}]).
                                   Defaults to None.
        """
        slide = self._add_slide()
        self._set_background(slide, self.colors['light_bg'])

        if badge_text:
            self._add_textbox(slide, badge_text, Inches(0.5), Inches(0.5), Inches(2), Inches(0.5), font_size=Pt(10), font_color=self.colors['grey_text2'])

        self._add_textbox(slide, icon, Inches(0.5), Inches(1.5), Inches(1.5), Inches(1.5), font_size=Pt(60), font_bold=True, font_color=self.colors['dark_text'], alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, name, Inches(2.25), Inches(1.5), Inches(7.25), Inches(0.75), font_size=Pt(36), font_bold=True, font_color=self.colors['dark_text'])
        if subtitle:
            self._add_textbox(slide, subtitle, Inches(2.25), Inches(2.25), Inches(7.25), Inches(0.5), font_size=Pt(18), font_color=self.colors['grey_text2'])

        self._add_textbox(slide, headline, Inches(0.5), Inches(3.25), Inches(9), Inches(0.5), font_size=Pt(18), font_color=self.colors['dark_text'])
        
        if description:
            self._add_textbox(slide, description, Inches(0.5), Inches(3.75), Inches(9), Inches(0.75), font_size=Pt(12), font_color=self.colors['grey_text3'])

        if stats:
            for i, stat in enumerate(stats):
                self._add_textbox(slide, stat['value'], Inches(0.5 + i*3), Inches(4.5), Inches(2.5), Inches(0.5), font_size=Pt(18), font_bold=True, font_color=self.colors['dark_text'], alignment=PP_ALIGN.CENTER)
                self._add_textbox(slide, stat['label'], Inches(0.5 + i*3), Inches(5.0), Inches(2.5), Inches(0.25), font_size=Pt(10), font_color=self.colors['grey_text2'], alignment=PP_ALIGN.CENTER)

    def add_grid_slide(self, cards, badge_text=None, subtitle=None):
        """
        Adds a grid-style slide with multiple cards.

        Args:
            cards (list): A list of dictionaries, each representing a card
                          (e.g., {"icon": "☀️", "title": "Card Title", "desc": "Description"}).
            badge_text (str, optional): Text for a badge/label. Defaults to None.
            subtitle (str, optional): A subtitle for the grid slide. Defaults to None.
        """
        slide = self._add_slide()
        self._set_background(slide, self.colors['dark_bg'])

        if badge_text:
            self._add_textbox(slide, badge_text, Inches(0.5), Inches(0.5), Inches(2), Inches(0.5), font_size=Pt(10), font_color=self.colors['grey_text1'])
        if subtitle:
            self._add_textbox(slide, subtitle, Inches(7.5), Inches(0.5), Inches(2), Inches(0.5), font_size=Pt(10), font_color=self.colors['grey_text1'], alignment=PP_ALIGN.RIGHT)

        for i, card in enumerate(cards):
            x = Inches(0.5 + i * 2.375)
            # Add a rectangular shape for the card background
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(2.125), Inches(3.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors['card_bg']
            shape.line.fill.background()

            self._add_textbox(slide, card['icon'], x, Inches(1.75), Inches(2.125), Inches(0.75), font_size=Pt(24), alignment=PP_ALIGN.CENTER, font_color=self.colors['white_text'])
            self._add_textbox(slide, card['title'], x, Inches(2.5), Inches(2.125), Inches(0.5), font_size=Pt(14), font_bold=True, alignment=PP_ALIGN.CENTER, font_color=self.colors['white_text'])
            self._add_textbox(slide, card['desc'], x, Inches(3.0), Inches(2.125), Inches(1.25), font_size=Pt(10), alignment=PP_ALIGN.CENTER, font_color=self.colors['grey_text2'])
            if 'label' in card:
                self._add_textbox(slide, card['label'], x, Inches(4.25), Inches(2.125), Inches(0.5), font_size=Pt(9), alignment=PP_ALIGN.CENTER, font_color=self.colors['grey_text3'])

    def add_features_slide(self, features, badge_text=None, subtitle=None, quote=None, quote_author=None):
        """
        Adds a slide showcasing features, with an optional quote.

        Args:
            features (list): A list of dictionaries, each with 'title' and 'desc'
                             for each feature.
            badge_text (str, optional): Text for a badge/label. Defaults to None.
            subtitle (str, optional): A subtitle for the features slide. Defaults to None.
            quote (str, optional): A quote to display on the slide. Defaults to None.
            quote_author (str, optional): The author of the quote. Defaults to None.
        """
        slide = self._add_slide()
        self._set_background(slide, self.colors['light_bg'])

        if badge_text:
            self._add_textbox(slide, badge_text, Inches(0.5), Inches(0.5), Inches(2), Inches(0.5), font_size=Pt(10), font_color=self.colors['grey_text2'])
        if subtitle:
            self._add_textbox(slide, subtitle, Inches(0.5), Inches(1.0), Inches(4), Inches(0.5), font_size=Pt(24), font_bold=True, font_color=self.colors['dark_text'])
        
        for i, feature in enumerate(features):
            y = Inches(1.75 + i * 0.8)
            self._add_textbox(slide, feature['title'], Inches(0.5), y, Inches(4), Inches(0.4), font_size=Pt(14), font_bold=True, font_color=self.colors['dark_text'])
            self._add_textbox(slide, feature['desc'], Inches(5.0), y, Inches(4.5), Inches(0.4), font_size=Pt(11), font_color=self.colors['grey_text3'])

        if quote:
            self._add_textbox(slide, f'"{quote}"', Inches(0.5), Inches(4.5), Inches(9), Inches(0.5), font_size=Pt(14), font_italic=True, font_color=self.colors['grey_text2'])
            if quote_author:
                 self._add_textbox(slide, f"- {quote_author}", Inches(7), Inches(5.0), Inches(2.5), Inches(0.25), font_size=Pt(11), font_color=self.colors['grey_text3'], alignment=PP_ALIGN.RIGHT)

    def add_tips_slide(self, tips, badge_text=None, subtitle=None):
        """
        Adds a slide displaying a list of tips or actions.

        Args:
            tips (list): A list of dictionaries, each with 'icon', 'title', and 'desc'
                         for each tip.
            badge_text (str, optional): Text for a badge/label. Defaults to None.
            subtitle (str, optional): A subtitle for the tips slide. Defaults to None.
        """
        slide = self._add_slide()
        self._set_background(slide, self.colors['dark_bg'])
        
        if badge_text:
            self._add_textbox(slide, badge_text, Inches(0.5), Inches(0.5), Inches(2), Inches(0.5), font_size=Pt(10), font_color=self.colors['grey_text1'])
        if subtitle:
            self._add_textbox(slide, subtitle, Inches(0.5), Inches(1.0), Inches(4), Inches(0.5), font_size=Pt(24), font_bold=True, font_color=self.colors['white_text'])

        for i, tip in enumerate(tips):
            x = Inches(0.5 + i * 3.16)
            self._add_textbox(slide, tip['icon'], x, Inches(2.0), Inches(3), Inches(0.75), font_size=Pt(28), font_color=self.colors['white_text'])
            self._add_textbox(slide, tip['title'], x, Inches(2.75), Inches(3), Inches(0.5), font_size=Pt(14), font_bold=True, font_color=self.colors['white_text'])
            self._add_textbox(slide, tip['desc'], x, Inches(3.25), Inches(3), Inches(1.5), font_size=Pt(11), font_color=self.colors['grey_text2'])

    def add_closing_slide(self, title, subtitle=None, badge_text=None, footer_items=None, closing_text=None):
        """
        Adds a closing slide to the presentation.

        Args:
            title (str): The main title for the closing slide.
            subtitle (str, optional): A subtitle. Defaults to None.
            badge_text (str, optional): Text for a badge/label. Defaults to None.
            footer_items (list, optional): A list of dictionaries for footer information.
                                          Defaults to None.
            closing_text (str, optional): Additional closing message. Defaults to None.
        """
        slide = self._add_slide()
        self._set_background(slide, self.colors['dark_bg'])
        
        if badge_text:
            self._add_textbox(slide, badge_text, Inches(0.5), Inches(0.5), Inches(2), Inches(0.5), font_size=Pt(10), font_color=self.colors['grey_text1'])

        self._add_textbox(slide, title, Inches(0.5), Inches(2.0), Inches(9), Inches(1.0), font_size=Pt(44), font_bold=True, font_color=self.colors['white_text'], alignment=PP_ALIGN.CENTER)
        if subtitle:
            self._add_textbox(slide, subtitle, Inches(0.5), Inches(3.0), Inches(9), Inches(0.5), font_size=Pt(16), font_color=self.colors['grey_text1'], alignment=PP_ALIGN.CENTER)
        
        if footer_items:
            for i, item in enumerate(footer_items):
                self._add_textbox(slide, f"{item['label']}: {item['value']}", Inches(0.5 + i*3), Inches(4.8), Inches(2.5), Inches(0.5), font_size=Pt(10), font_color=self.colors['grey_text3'])
        
        if closing_text:
            self._add_textbox(slide, closing_text, Inches(7), Inches(4.8), Inches(2.5), Inches(0.5), font_size=Pt(18), font_color=self.colors['grey_text1'], alignment=PP_ALIGN.RIGHT)

    def _add_textbox(self, slide, text, left, top, width, height, font_size=Pt(12), font_bold=False, font_italic=False, font_color=RGBColor(0,0,0), alignment=PP_ALIGN.LEFT):
        """
        Helper method to add a formatted textbox to a slide.

        Args:
            slide (Slide): The slide object to add the textbox to.
            text (str): The text content of the textbox.
            left (Inches): The x-coordinate of the left edge of the textbox.
            top (Inches): The y-coordinate of the top edge of the textbox.
            width (Inches): The width of the textbox.
            height (Inches): The height of the textbox.
            font_size (Pt, optional): The font size. Defaults to Pt(12).
            font_bold (bool, optional): Whether the font should be bold. Defaults to False.
            font_italic (bool, optional): Whether the font should be italic. Defaults to False.
            font_color (RGBColor, optional): The font color. Defaults to RGBColor(0,0,0).
            alignment (PP_ALIGN, optional): The text alignment. Defaults to PP_ALIGN.LEFT.
        """
        textbox_shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox_shape.text_frame
        text_frame.clear() 
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.name = 'Arial'
        paragraph.font.size = font_size
        paragraph.font.bold = font_bold
        paragraph.font.italic = font_italic
        paragraph.font.color.rgb = font_color
        paragraph.alignment = alignment
        text_frame.word_wrap = True

    def save(self, filename):
        """
        Saves the presentation to a specified file.

        Args:
            filename (str): The path and name of the file to save the presentation as.
        """
        self.presentation.save(filename)
